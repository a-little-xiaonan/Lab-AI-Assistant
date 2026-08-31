"""文档加载器测试：md/txt 直读、GBK 回退、兜底链（无需 API key）。

PDF/DOCX/XLSX/PPTX 用库现场生成（pymupdf/python-docx/openpyxl/python-pptx 已装）。
"""
import tempfile
from pathlib import Path

import pytest

from app.core import document_loader
from app.core.document_loader import DocumentParseError, UnsupportedFormatError

FIXTURES = Path(__file__).parent / "fixtures"


# ---------- 直读格式 ----------

def test_md_heading_sections():
    fmt, elements = document_loader.load(FIXTURES / "sample.md")
    assert fmt == "md"
    assert len(elements) >= 5  # 每个标题一个结构块
    assert any(e.text.startswith("# 千问大模型产品手册") for e in elements)


def test_txt_paragraphs():
    fmt, elements = document_loader.load(FIXTURES / "sample.txt")
    assert fmt == "txt"
    assert len(elements) == 3


def test_gbk_text_fallback():
    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "gbk.txt"
        p.write_bytes("这是中文测试。\n\n第二段。".encode("gbk"))
        fmt, elements = document_loader.load(p)
    assert fmt == "txt"
    assert "这是中文测试。" in elements[0].text


# ---------- 兜底链 ----------

def test_unknown_text_ext_goes_through_fallback():
    # .log 不在白名单：unstructured 自动分区（L2）能解析 → 文本正常入库
    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "app.log"
        p.write_text("2026-08-30 INFO 服务启动成功\n2026-08-30 ERROR 连接失败", encoding="utf-8")
        fmt, elements = document_loader.load(p)
    assert fmt == "auto"
    assert "服务启动成功" in " ".join(e.text for e in elements)


def test_l3_text_probe_when_unstructured_fails(monkeypatch):
    # 模拟 unstructured 不可用 → 降级 L3 文本试探（UTF-8 解码）
    def boom(*_a, **_k):
        raise RuntimeError("unstructured down")

    monkeypatch.setattr("unstructured.partition.auto.partition", boom)
    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "app.log"
        p.write_text("2026-08-30 INFO 服务启动成功\n\n2026-08-30 ERROR 连接失败", encoding="utf-8")
        fmt, elements = document_loader.load(p)
    assert fmt == "txt"
    assert len(elements) == 2  # 按段落切


def test_unknown_binary_rejected_gracefully():
    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "data.bin"
        p.write_bytes(bytes(range(256)) * 8)  # 不可解码的二进制
        with pytest.raises(UnsupportedFormatError):
            document_loader.load(p)


def test_detect_format():
    assert document_loader.detect_format(Path("a.PDF")) == "pdf"  # 大小写不敏感
    assert document_loader.detect_format(Path("a.xyz")) is None


# ---------- 现场生成的格式 ----------

def test_pdf_pages_and_cross_page_merge():
    import pymupdf

    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "demo.pdf"
        with pymupdf.open() as doc:
            # 第 1 页：一个未以句末标点结尾的段落（会与第 2 页首块合并）
            page1 = doc.new_page()
            page1.insert_text((72, 72), "这是一个跨页段落，", fontsize=12, fontname="china-s")
            # 第 2 页：续写该段落 + 独立段落
            page2 = doc.new_page()
            page2.insert_text((72, 72), "后半部分内容。", fontsize=12, fontname="china-s")
            page2.insert_text((72, 110), "独立段落的内容。", fontsize=12, fontname="china-s")
            doc.save(p)
        fmt, elements = document_loader.load(p)
    assert fmt == "pdf"
    assert any("跨页段落" in e.text and "后半部分" in e.text for e in elements)  # 跨页合并
    assert any(e.text == "独立段落的内容。" and e.page == 1 for e in elements)  # 第 2 页 → page=1 (0-based)


def test_docx_heading_sections():
    from docx import Document as Docx

    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "demo.docx"
        doc = Docx()
        doc.add_heading("第一章 产品介绍", level=1)
        doc.add_paragraph("千问大模型支持多轮对话。")
        doc.add_heading("第二章 定价", level=1)
        doc.add_paragraph("按 token 计费。")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "模型"
        table.cell(0, 1).text = "价格"
        table.cell(1, 0).text = "qwen-plus"
        table.cell(1, 1).text = "0.0008"
        doc.save(p)
        fmt, elements = document_loader.load(p)
    assert fmt == "docx"
    sections = [e for e in elements if "第一章" in e.text]
    assert len(sections) == 1
    assert "多轮对话" in sections[0].text  # 章节内容并入标题节
    assert any("模型 | 价格" in e.text for e in elements)  # 表格


def test_xlsx_sheet_blocks_with_header_repeat():
    import openpyxl

    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "demo.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "价格"
        ws.append(["模型", "价格"])
        for i in range(60):
            ws.append([f"qwen-{i}", f"{i}.001"])
        wb.save(p)
        fmt, elements = document_loader.load(p)
    assert fmt == "xlsx"
    assert len(elements) == 2  # ~50 行/块 → 61 行 = 2 块
    assert elements[0].sheet_name == "价格"
    assert elements[0].row_range == "1-50"
    assert elements[1].row_range == "51-61"
    assert elements[1].text.startswith("模型 | 价格")  # 表头重复


def test_pptx_grouped_by_slide():
    from pptx import Presentation

    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "demo.pptx"
        prs = Presentation()
        for title, body in [("第一页", "这是第一页的正文。"), ("第二页", "这是第二页的正文。")]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
        prs.save(p)
        fmt, elements = document_loader.load(p)
    assert fmt == "pptx"
    assert len(elements) == 2
    assert elements[0].slide_number == 1
    assert "第一页" in elements[0].text
    assert elements[1].slide_number == 2
    assert "第二页" in elements[1].text


def test_html_elements():
    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "demo.html"
        p.write_text("<html><body><h1>标题</h1><p>第一段。</p><p>第二段。</p></body></html>", encoding="utf-8")
        fmt, elements = document_loader.load(p)
    assert fmt == "html"
    assert any(e.text == "第一段。" for e in elements)
